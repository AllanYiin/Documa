"""PPTX parser adapter."""

from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Any

from documa.adapters.base import ParseOptions, ParserAdapter
from documa.core.errors import DocumaError, DocumaErrorDetail
from documa.core.ir import BlockIR, BlockType, Confidence, DocumentIR, PageIR, TableIR, TextContent


def _load_pptx():
    try:
        from pptx import Presentation  # type: ignore

        return Presentation
    except ImportError as exc:
        raise DocumaError(
            DocumaErrorDetail(
                code="PPTX_DEPENDENCY_NOT_INSTALLED",
                message="python-pptx is required for PptxAdapter.",
                recoverable=True,
                suggested_action="Install the optional dependency: pip install documa[pptx]",
            )
        ) from exc


def _document_id(source_path: Path, size: int) -> str:
    digest = hashlib.sha256(f"{source_path.resolve()}\n{size}".encode("utf-8")).hexdigest()[:16]
    return f"doc_pptx_{digest}"


def _points(value: Any) -> float:
    return float(getattr(value, "pt", value or 0.0))


def _shape_bbox(shape: Any) -> tuple[float, float, float, float] | None:
    try:
        left = _points(shape.left)
        top = _points(shape.top)
        width = _points(shape.width)
        height = _points(shape.height)
    except Exception:
        return None
    return (left, top, left + width, top + height)


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    paragraphs = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs).strip()


def _is_title_shape(slide: Any, shape: Any) -> bool:
    if str(getattr(shape, "name", "") or "").casefold().startswith("title"):
        return True
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        placeholder_type = str(shape.placeholder_format.type).casefold()
    except Exception:
        return False
    return "title" in placeholder_type


def _paragraph_levels(shape: Any) -> list[int]:
    if not getattr(shape, "has_text_frame", False):
        return []
    return [int(getattr(paragraph, "level", 0) or 0) for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()]


def _table_rows(shape: Any) -> list[list[str | None]]:
    if not getattr(shape, "has_table", False):
        return []
    rows: list[list[str | None]] = []
    for row in shape.table.rows:
        cells = [cell.text.strip() or None for cell in row.cells]
        if any(cells):
            rows.append(cells)
    return rows


def _table_text(rows: list[list[str | None]]) -> str:
    return "\n".join(" | ".join("" if cell is None else cell for cell in row) for row in rows)


class PptxAdapter(ParserAdapter):
    """Parse PPTX slides into slide/page-oriented Documa IR."""

    name = "pptx"

    def parse(self, source: str | Path, options: ParseOptions | None = None) -> DocumentIR:
        options = options or ParseOptions()
        source_path = Path(source)
        Presentation = _load_pptx()

        try:
            presentation = Presentation(str(source_path))
        except Exception as exc:
            raise DocumaError(
                DocumaErrorDetail(
                    code="PPTX_OPEN_FAILED",
                    message=f"Unable to open PPTX: {source_path}",
                    recoverable=True,
                    suggested_action="Check whether the file exists and is a valid .pptx file.",
                    context={"source": str(source_path), "error": str(exc)},
                )
            ) from exc

        document = DocumentIR(
            id=_document_id(source_path, source_path.stat().st_size if source_path.exists() else 0),
            source_name=str(source_path),
            parser=self.name,
            metadata={
                "adapter": self.name,
                "format": "pptx",
                "languages": list(options.languages),
                "slide_count": len(presentation.slides),
            },
        )

        table_index = 0
        for slide_index, slide in enumerate(presentation.slides, start=1):
            page = PageIR(
                id=f"slide_{slide_index}",
                page_number=slide_index,
                width=_points(presentation.slide_width),
                height=_points(presentation.slide_height),
                metadata={"source": "pptx_slide", "slide_index": slide_index},
            )
            ordered_shapes = sorted(
                list(slide.shapes),
                key=lambda shape: (_points(getattr(shape, "top", 0)), _points(getattr(shape, "left", 0))),
            )
            for order, shape in enumerate(ordered_shapes, start=1):
                rows = _table_rows(shape)
                if rows:
                    table_index += 1
                    block = BlockIR(
                        id=f"pptx_s{slide_index}_table{table_index:04d}",
                        type=BlockType.TABLE,
                        page_number=slide_index,
                        text=TextContent(_table_text(rows)),
                        bbox=_shape_bbox(shape),
                        confidence=Confidence.MEDIUM,
                        order_index=order,
                        source_refs=[f"pptx:slide:{slide_index}:table:{table_index}"],
                        metadata={
                            "source_type": "pptx_table",
                            "table_rows": rows,
                            "table_title": rows[0][0] if rows and rows[0] else None,
                        },
                    )
                    page.blocks.append(block)
                    document.tables.append(
                        TableIR(
                            id=f"table_{block.id}",
                            block_id=block.id,
                            rows=rows,
                            confidence=Confidence.MEDIUM,
                            metadata={"source": "pptx_table"},
                        )
                    )
                    continue

                text = _shape_text(shape)
                if not text:
                    continue
                is_title = _is_title_shape(slide, shape)
                metadata = {
                    "source_type": "pptx_text_shape",
                    "shape_name": getattr(shape, "name", None),
                    "paragraph_levels": _paragraph_levels(shape),
                    "paragraph_title": None if is_title else text.splitlines()[0],
                }
                if is_title:
                    metadata["heading_level"] = 1
                page.blocks.append(
                    BlockIR(
                        id=f"pptx_s{slide_index}_b{order:04d}",
                        type=BlockType.HEADING if is_title else BlockType.PARAGRAPH,
                        page_number=slide_index,
                        text=TextContent(text),
                        bbox=_shape_bbox(shape),
                        confidence=Confidence.HIGH,
                        order_index=order,
                        source_refs=[f"pptx:slide:{slide_index}:shape:{order}"],
                        metadata=metadata,
                    )
                )
            document.pages.append(page)

        return document
