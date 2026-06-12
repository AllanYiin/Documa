import base64
from email.message import EmailMessage
import json
import sys
import tempfile
import types
import unittest
from pathlib import Path

from documa.cli import main
from documa.core.ir import BlockType
from documa.interfaces import call_documa_tool


class DocumentFormatAdapterTests(unittest.TestCase):
    def test_docx_process_builds_sections_tables_and_searchable_blocks(self):
        try:
            from docx import Document  # type: ignore
        except ImportError:
            self.skipTest("python-docx is required")

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            docx_path = tmp_path / "strategy.docx"
            out_dir = tmp_path / "out"

            document = Document()
            document.add_heading("策略總覽", level=1)
            document.add_paragraph("Documa 會保留文件結構、人類閱讀順序與分塊閱讀線索。")
            table = document.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "格式"
            table.cell(0, 1).text = "支援"
            table.cell(1, 0).text = "DOCX"
            table.cell(1, 1).text = "段落與表格"
            document.save(docx_path)

            result = call_documa_tool(
                "documa_process",
                {"source": str(docx_path), "out": str(out_dir), "export_formats": ["block-json"]},
            )

            self.assertFalse(result["isError"])
            payload = result["structuredContent"]
            self.assertEqual(payload["parser"], "docx")
            self.assertEqual(payload["status"], "ok")

            ir = json.loads((out_dir / "documa.ir.json").read_text(encoding="utf-8"))
            block_types = [block["type"] for page in ir["pages"] for block in page["blocks"]]
            self.assertIn(BlockType.HEADING.value, block_types)
            self.assertIn(BlockType.TABLE.value, block_types)

            search = call_documa_tool(
                "documa_search_blocks",
                {"ir_path": str(out_dir / "documa.ir.json"), "query": "分塊閱讀", "limit": 5},
            )
            self.assertFalse(search["isError"])
            self.assertGreaterEqual(len(search["structuredContent"]["results"]), 1)

    def test_pptx_process_treats_slides_as_pages(self):
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            self.skipTest("python-pptx is required")

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            pptx_path = tmp_path / "briefing.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "產品路線"
            body = slide.shapes.placeholders[1].text_frame
            body.text = "支援 PDF、Markdown、DOCX、PPTX、HTML"
            rows = cols = 2
            table = slide.shapes.add_table(rows, cols, 0, 0, presentation.slide_width // 2, presentation.slide_height // 4).table
            table.cell(0, 0).text = "項目"
            table.cell(0, 1).text = "狀態"
            table.cell(1, 0).text = "PPTX"
            table.cell(1, 1).text = "已解析"
            presentation.save(pptx_path)

            result = call_documa_tool("documa_process", {"source": str(pptx_path)})
            document = result["structuredContent"]["document"]

            self.assertFalse(result["isError"])
            self.assertEqual(result["structuredContent"]["parser"], "pptx")
            self.assertEqual(result["structuredContent"]["page_count"], 1)
            self.assertTrue(any(block["type"] == "section" and block["title"] == "產品路線" for block in document["document_blocks"]))
            self.assertTrue(any(block["type"] == "table" for block in document["document_blocks"]))

    def test_html_cli_process_preserves_dom_order_headings_tables_and_links(self):
        html = """<!doctype html>
<html lang="zh-Hant">
  <head><title>Documa 文件</title></head>
  <body>
    <section>
      <h1>文件理解</h1>
      <p>HTML 解析應保留 <a href="https://example.com/ref">來源連結</a> 與閱讀順序。</p>
      <table>
        <tr><th>格式</th><th>功能</th></tr>
        <tr><td>HTML</td><td>DOM order</td></tr>
      </table>
    </section>
  </body>
</html>
"""
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            html_path = tmp_path / "sample.html"
            out_dir = tmp_path / "out"
            html_path.write_text(html, encoding="utf-8")

            from io import StringIO

            old_stdout = sys.stdout
            sys.stdout = StringIO()
            try:
                exit_code = main(["process", str(html_path), "--out", str(out_dir), "--export-format", "block-json"])
                output = json.loads(sys.stdout.getvalue())
            finally:
                sys.stdout = old_stdout

            self.assertEqual(exit_code, 0)
            self.assertEqual(output["parser"], "html")
            ir = json.loads((out_dir / "documa.ir.json").read_text(encoding="utf-8"))
            blocks = [block for page in ir["pages"] for block in page["blocks"]]
            self.assertEqual([block["type"] for block in blocks], ["heading", "paragraph", "table"])
            self.assertEqual(blocks[1]["metadata"]["links"][0]["href"], "https://example.com/ref")

    def test_eml_process_extracts_headers_body_and_attachments(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            eml_path = tmp_path / "message.eml"
            out_dir = tmp_path / "out"

            message = EmailMessage()
            message["Subject"] = "專案更新"
            message["From"] = "Alice <alice@example.com>"
            message["To"] = "Bob <bob@example.com>"
            message["Cc"] = "Carol <carol@example.com>"
            message["Date"] = "Fri, 12 Jun 2026 09:30:00 +0800"
            message.set_content("請閱讀附件中的規格，並回覆主要風險。")
            message.add_attachment(b"attachment bytes", maintype="application", subtype="octet-stream", filename="spec.txt")
            eml_path.write_bytes(message.as_bytes())

            result = call_documa_tool(
                "documa_process",
                {"source": str(eml_path), "out": str(out_dir), "export_formats": ["block-json"]},
            )

            self.assertFalse(result["isError"])
            self.assertEqual(result["structuredContent"]["parser"], "eml")
            ir = json.loads((out_dir / "documa.ir.json").read_text(encoding="utf-8"))
            email_metadata = ir["metadata"]["email"]
            self.assertEqual(email_metadata["subject"], "專案更新")
            self.assertEqual(email_metadata["sender"], "Alice <alice@example.com>")
            self.assertEqual(email_metadata["receiver"], ["Bob <bob@example.com>"])
            self.assertEqual(email_metadata["cc"], ["Carol <carol@example.com>"])
            self.assertEqual(email_metadata["date_iso"], "2026-06-12T09:30:00+08:00")
            self.assertEqual(email_metadata["attachments"][0]["filename"], "spec.txt")
            self.assertTrue((out_dir / "assets" / email_metadata["attachments"][0]["asset_ref"]).exists())

            blocks = [block for page in ir["pages"] for block in page["blocks"]]
            self.assertTrue(any(block["id"] == "email_body" and "主要風險" in block["text"]["raw_text"] for block in blocks))

            search = call_documa_tool(
                "documa_search_blocks",
                {"ir_path": str(out_dir / "documa.ir.json"), "query": "主要風險", "limit": 5},
            )
            self.assertFalse(search["isError"])
            self.assertGreaterEqual(len(search["structuredContent"]["results"]), 1)

    def test_msg_process_uses_extract_msg_boundary(self):
        class FakeAttachment:
            longFilename = "outlook-note.txt"
            mimetype = "text/plain"
            cid = "cid-1"
            data = b"msg attachment"

        class FakeMessage:
            subject = "Outlook 會議紀錄"
            sender = "Dana <dana@example.com>"
            to = "Eli <eli@example.com>; Fran <fran@example.com>"
            cc = "Grace <grace@example.com>"
            bcc = ""
            date = "2026-06-12 10:00:00+08:00"
            body = "這封 Outlook MSG 包含行動項目與附件。"
            attachments = [FakeAttachment()]

            def close(self):
                self.closed = True

        previous = sys.modules.get("extract_msg")
        sys.modules["extract_msg"] = types.SimpleNamespace(openMsg=lambda _: FakeMessage())
        try:
            with tempfile.TemporaryDirectory() as tmp:
                tmp_path = Path(tmp)
                msg_path = tmp_path / "meeting.msg"
                out_dir = tmp_path / "out"
                msg_path.write_bytes(b"fake msg bytes")

                result = call_documa_tool(
                    "documa_process",
                    {"source": str(msg_path), "out": str(out_dir), "export_formats": ["block-json"]},
                )

                self.assertFalse(result["isError"])
                self.assertEqual(result["structuredContent"]["parser"], "msg")
                ir = json.loads((out_dir / "documa.ir.json").read_text(encoding="utf-8"))
                email_metadata = ir["metadata"]["email"]
                self.assertEqual(email_metadata["subject"], "Outlook 會議紀錄")
                self.assertEqual(email_metadata["sender"], "Dana <dana@example.com>")
                self.assertEqual(email_metadata["receiver"], ["Eli <eli@example.com>", "Fran <fran@example.com>"])
                self.assertEqual(email_metadata["attachments"][0]["filename"], "outlook-note.txt")
                self.assertTrue((out_dir / "assets" / email_metadata["attachments"][0]["asset_ref"]).exists())
                blocks = [block for page in ir["pages"] for block in page["blocks"]]
                self.assertTrue(any(block["id"] == "email_body" and "行動項目" in block["text"]["raw_text"] for block in blocks))
        finally:
            if previous is None:
                sys.modules.pop("extract_msg", None)
            else:
                sys.modules["extract_msg"] = previous

    def test_ipynb_process_preserves_cell_order_and_attachments(self):
        try:
            import nbformat  # type: ignore
        except ImportError:
            self.skipTest("nbformat is required")

        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            ipynb_path = tmp_path / "analysis.ipynb"
            out_dir = tmp_path / "out"
            notebook = nbformat.v4.new_notebook(
                metadata={
                    "kernelspec": {"name": "python3", "display_name": "Python 3"},
                    "language_info": {"name": "python"},
                }
            )
            markdown_cell = nbformat.v4.new_markdown_cell("# 分析摘要\n這份 notebook 會保留 cell order。")
            markdown_cell["attachments"] = {
                "chart.png": {"image/png": base64.b64encode(b"fake png").decode("ascii")},
            }
            code_cell = nbformat.v4.new_code_cell(
                "risk_score = 0.42\nprint(risk_score)",
                execution_count=1,
                outputs=[nbformat.v4.new_output("stream", name="stdout", text="0.42\n")],
            )
            notebook.cells = [markdown_cell, code_cell]
            with ipynb_path.open("w", encoding="utf-8", newline="\n") as handle:
                nbformat.write(notebook, handle)

            result = call_documa_tool(
                "documa_process",
                {"source": str(ipynb_path), "out": str(out_dir), "export_formats": ["block-json"]},
            )

            self.assertFalse(result["isError"])
            self.assertEqual(result["structuredContent"]["parser"], "ipynb")
            ir = json.loads((out_dir / "documa.ir.json").read_text(encoding="utf-8"))
            self.assertEqual(ir["metadata"]["cell_count"], 2)
            self.assertEqual(ir["metadata"]["kernelspec"]["name"], "python3")
            self.assertEqual(ir["metadata"]["attachments"][0]["filename"], "chart.png")
            self.assertTrue((out_dir / "assets" / ir["metadata"]["attachments"][0]["asset_ref"]).exists())
            blocks = [block for page in ir["pages"] for block in page["blocks"]]
            self.assertEqual([block["metadata"]["cell_type"] for block in blocks], ["markdown", "code"])
            self.assertEqual(blocks[0]["type"], BlockType.HEADING.value)
            self.assertIn("risk_score", blocks[1]["text"]["raw_text"])


if __name__ == "__main__":
    unittest.main()
